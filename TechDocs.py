# -*- coding: utf-8 -*-

#***************************************************************************
#*                                                                         *
#*   Copyright (c) 2019                                                    *
#*   This file is part of the FreeCAD CAx development system.              *
#*                                                                         *
#*   Authors:                                                              *
#*   Benjamin Alterauge (gift)                                             *
#*                                                                         *
#*   This program is free software; you can redistribute it and/or modify  *
#*   it under the terms of the GNU Lesser General Public License (LGPL)    *
#*   as published by the Free Software Foundation; either version 2 of     *
#*   the License, or (at your option) any later version.                   *
#*   for detail see the LICENCE text file.                                 *
#*                                                                         *
#*   This program is distributed in the hope that it will be useful,       *
#*   but WITHOUT ANY WARRANTY; without even the implied warranty of        *
#*   MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the         *
#*   GNU Library General Public License for more details.                  *
#*                                                                         *
#*   You should have received a copy of the GNU Library General Public     *
#*   License along with this program; if not, write to the Free Software   *
#*   Foundation, Inc., 59 Temple Place, Suite 330, Boston, MA  02111-1307  *
#*   USA                                                                   *
#*                                                                         *
#***************************************************************************

import FreeCAD, FreeCADGui, TechDocsTools
from pivy import coin

def makePresentationPlane():
    '''makePresentationPlane()'''
    obj=FreeCAD.ActiveDocument.addObject('Part::FeaturePython','PresentationPlane')
    PresentationPlane(obj)
    ViewProviderPresentationPlane(obj.ViewObject)
    sel = FreeCADGui.Selection
    obj.ViewObject.Transparency = 80
    obj.ViewObject.ShapeColor   = (0.0, 0.0, 0.8)
    obj.Objects = sel.getSelection()
    return obj

class PresentationPlane:
    '''Text...'''

    pageSize = {
        'SizeA4Paper' : [297.0, 210.0],
        'SizeA3Paper' : [420.0, 297.0],
    }

    def __init__(self, obj):
        self.Object = obj
        self.Type = 'PresentationPlane'
        obj.Proxy = self
        obj.addProperty('App::PropertyLinkList','Objects','Plane', 'The objects included in this selector.')
        obj.addProperty('App::PropertyEnumeration','PageSize','Plane', 'Choose the size of the page.')
        obj.addProperty('App::PropertyFloat','Scale','Plane', 'Scale factor.')
        obj.addProperty('App::PropertyLength','Discretize','Plane', 'Discretize length.').Discretize = 2.0
        obj.addProperty('App::PropertyFile', 'File', 'Exporter', 'File').File = ''
        obj.addProperty('App::PropertyEnumeration','Gateway','Exporter', 'Gateway')
        obj.Gateway = ['PPTX']
        setup = ['SizeA4Paper','SizeA3Paper']
        obj.PageSize = setup
        obj.PageSize = setup.index(setup[0])
        obj.Scale = 1.0

    def execute(self, obj):
        import Part, math
        w, h = PresentationPlane.pageSize[obj.PageSize]
        s = obj.Scale
        p = obj.Placement
        face = Part.makePlane(w / s, h / s, FreeCAD.Vector(w / s / -2.0, h / s / -2.0, 0))
        face.translate(p.Base)
        face.rotate(p.Base, p.Rotation.Axis, p.Rotation.Angle * 180 /  math.pi )
        v = FreeCAD.Vector(p.Base.x *-2.0, p.Base.y *-2.0, p.Base.z *-2.0)
        result = []
        for objs in obj.Objects:
            for edge in objs.Shape.Edges:
                tmp = face.makeParallelProjection(edge, obj.Placement.Rotation.Axis)
                tmp.translate(v)
                result.append(tmp)

        obj.Shape = Part.makeCompound(result)

    def getPageSize(self, obj):
        return PresentationPlane.pageSize[obj.PageSize]

    def LinesToPoints(self, obj):
        import Part
        pts = []
        for edg in obj.Shape.Edges:
            if type(edg.Curve) == Part.Line:
                pts.append([[edg.Vertexes[0].X, edg.Vertexes[0].Y, edg.Vertexes[0].Z],[edg.Vertexes[1].X, edg.Vertexes[1].Y, edg.Vertexes[1].Z]])
        return pts

    def ArcsToPoints(self, obj):
        import Part
        pts = []
        for edg in obj.Shape.Edges:
            if type(edg.Curve) == Part.Circle:
                ptg = []
                for pt in edg.Curve.discretize(Distance=obj.Discretize,First=edg.FirstParameter,Last=edg.LastParameter):
                    ptg.append([pt.x, pt.y, pt.z])

                pts.append(ptg)
        return pts
        
    def exportNow(self, obj):
        if (obj.Gateway == 'PPTX') and not (obj.File == ''):
            from pptx.enum.shapes import MSO_CONNECTOR
            from pptx.util import Mm
            from pptx import Presentation

            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            w, h = self.getPageSize(obj)
            w = w / 2
            h = h / 2
            for ln in self.LinesToPoints(obj):
                line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Mm(ln[0][0]+w), Mm((ln[0][1]*-1.0)+h), Mm(ln[1][0]+w), Mm((ln[1][1]*-1.0)+h) )

            prs.save(obj.File)

    def bundlePoints(self, obj):
        result = self.LinesToPoints(obj)
        result.extend(self.ArcsToPoints(obj))
        return result

    def __getstate__(self):
        return None

    def __setstate__(self, state):
        return None

class CoinTemplate:

    def __init__(self, spGrp):
        self.__color     = coin.SoBaseColor()
        self.__drawStyle = coin.SoDrawStyle()
        spGrp.addChild(self.__color)
        spGrp.addChild(self.__drawStyle)
        self.__group = spGrp
        self.style()

    def style(self, sd={'color': (0.,0.,0.), 'pointSize': 1.0, 'lineWidth': 1.0} ):
        if 'color' in sd:
            self.__color.rgb.setValue(sd['color'])
        if 'pointSize' in sd:
            self.__drawStyle.pointSize = sd['pointSize']
        if 'lineWidth' in sd:
            self.__drawStyle.lineWidth = sd['lineWidth']
        if 'style' in sd:
            self.__drawStyle.style = sd['style']

class CoinCoordinate():

    def __init__(self, spGrp):
        self.__points = coin.SoCoordinate3()
        spGrp.addChild(self.__points)
        self.__group = spGrp

    def getPoints(self):
        return [p.getValue() for p in self.__points.point.getValues()]

    def setPoints(self, pts):
        num = len(pts)
        self.__points.point.setValue(0,0,0)
        self.__points.point.setValues(0, num, pts)

    def push(self, pt):
        pts = self.GetPoints()
        pts.append(pt)
        self.setPoints(pts)

    def pop(self, i=1):
        pts = self.GetPoints()
        r = pts.pop(i)
        self.setPoints(pts)
        return r

class CoinPolygon(CoinTemplate):

    def __init__(self, spGrp):
        CoinTemplate.__init__(self, spGrp)
        self.__lines = coin.SoIndexedLineSet()
        spGrp.addChild(self.__lines)
        self.style({'style': coin.SoDrawStyle.LINES})

    def coordInit(self):
        self.__lines.coordIndex.setValue(0)

    def push(self, num, pt):
        self.__lines.coordIndex.set1Value(num, pt)

class CoinEdgeSet(CoinTemplate):

    def __init__(self, spGrp):
        CoinTemplate.__init__(self, spGrp)
        self.__lineset = coin.SoType.fromName('SoBrepEdgeSet').createInstance()
        self.__lineset.highlightIndex = -1
        self.__lineset.selectionIndex = 0
        spGrp.addChild(self.__lineset)

    def coordInit(self):
        self.__lineset.coordIndex.setValue(0)

    def push(self, num, pt):
        self.__lineset.coordIndex.set1Value(num, pt)

class CoinPoints(CoinTemplate):

    def __init__(self, spGrp):
        CoinTemplate.__init__(self, spGrp)
        self.__pointset = coin.SoType.fromName('SoBrepPointSet').createInstance()
        spGrp.addChild(self.__pointset)
        self.style({'style': coin.SoDrawStyle.POINTS})

    def setNumPoints(self, num):
        self.__pointset.numPoints.setValue(num)

    def incNumPoints(self, num):
        n = self.__pointset.numPoints.getValue()
        self.__pointset.numPoints.setValue(n+num)

class CoinMarkers(CoinTemplate):

    def __init__(self, spGrp):
        CoinTemplate.__init__(self, spGrp)
        self.__markers = coin.SoMarkerSet()
        spGrp.addChild(self.__markers)
        self.__markers.startIndex  = 0
        self.__markers.numPoints   = 1
        self.__markers.markerIndex = coin.SoMarkerSet.CIRCLE_LINE_9_9

    def setStart(self, i):
        self.__markers.startIndex = i

    def setNum(self, n):
        self.__markers.numPoints = n

class CoinGroup():

    def __init__(self):
        self.group = coin.SoGroup()
        self.__scale = coin.SoScale()
        self.group.addChild(self.__scale)

    def addCoinObject(self, cls, name):
        exec('self.'+name+' = '+cls+'(self.group)')

    def addCoordinate(self, name='coords'):
        self.addCoinObject('CoinCoordinate', name)

    def addPolygon(self, name='polygon'):
        self.addCoinObject('CoinPolygon', name)

    def addPoints(self, name='points'):
        self.addCoinObject('CoinPoints', name)

    def addMarkers(self, name='markers'):
        self.addCoinObject('CoinMarkers', name)

    def addEdgeSet(self, name='edgeset'):
        self.addCoinObject('CoinEdgeSet', name)

class ViewProviderTemplate:

    def __init__(self, vobj):
        vobj.Proxy = self
        self.Object = vobj.Object

    def __getstate__(self):
        return None

    def __setstate__(self, state):
        return None

    def getProxyType(self, obj):
        if hasattr(obj, 'Proxy'):
           if hasattr(obj.Proxy, 'Type'):
              return obj.Proxy.Type
        return None

class ViewProviderPresentationPlane(ViewProviderTemplate):

    def __init__(self, vobj):
        ViewProviderTemplate.__init__(self, vobj)
        self.Restore = None

    def getIcon(self):
        return TechDocsTools.Settings.icon('PresentationPlane.svg')


    def attach(self, obj):
        self.Object = obj.Object

#    def claimChildren(self):
#        return self.Object.Objects
        
    def doubleClicked(self, obj):
        self.Object.Proxy.exportNow(self.Object)

    def getDisplayModes(self,obj):
        return ['Presentation']

    def getDefaultDisplayMode(self):
        return u'Presentation'

    def attach(self, obj):
        self.Object = obj.Object
        self.presentation = CoinGroup()
        self.presentation.addCoordinate()
        self.presentation.addPoints()
        self.presentation.addPolygon()
        self.presentation.points.style({'color': (1.0, 0., 1.0), 'pointSize': 4.0})
        self.presentation.polygon.style({'color': (1.0, 0., 1.0), 'lineWidth': 1.5})

        self.presentation.addCoordinate('pt_coords')
        self.presentation.addPoints('pt_points')
        self.presentation.addCoordinate('ed_coords')
        self.presentation.addEdgeSet()
        self.presentation.pt_points.style(
             {'color': (1., 1., 1.), 'pointSize': 2.0})
        self.presentation.edgeset.style(
             {'color': (0., 0., 0.), 'lineWidth': 1.0})

        obj.addDisplayMode(self.presentation.group, u'Presentation')



    def updateData(self, fp, prop):
        if prop == 'Shape':
            import Part, math
            pre = self.presentation
            h, w = fp.Proxy.getPageSize(fp)
            s = fp.Scale

            points = [[h / -2.0 / s, w / 2.0 / s, 0.0], [h / -2.0 / s, w / -2.0 / s, 0.0],
                      [h / 2.0 / s, w / -2.0 / s, 0.0], [h / 2.0  / s, w / 2.0 / s, 0.0]]

            pre.points.setNumPoints(len(points))
            pre.coords.setPoints(points)
            pre.polygon.coordInit()

            for cnt in range(4):
                pre.polygon.push(cnt, cnt)

            pre.polygon.push(4, 0)
            pre.polygon.push(5, -1)

            ptg = fp.Proxy.bundlePoints(fp)

            pts = []
            for pt in ptg:
                if len(pt) > 1:
                    pts.append(pt[0])
                    pts.append(pt[-1])

            pre.pt_points.setNumPoints(len(pts))
            pre.pt_coords.setPoints(pts)

            pts = [b for a in ptg for b in a]
            pre.ed_coords.setPoints(pts)
            pre.edgeset.coordInit()

            cnt = 0
            ptn = 0
            for pt in ptg:
                if len(pt) > 1:
                    for i in pt:
                        pre.edgeset.push(cnt, ptn)
                        cnt += 1
                        ptn += 1

                    pre.edgeset.push(cnt, -1)
                    cnt += 1



